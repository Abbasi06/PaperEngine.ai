import torch
import json
import random
import os
import warnings
from transformers import AutoModelForCausalLM, AutoTokenizer, BitsAndBytesConfig
from peft import PeftModel
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Suppress warnings during model loading for cleaner output
warnings.filterwarnings("ignore", category=FutureWarning)

# --- Configuration Constants ---
MODEL_NAME = "microsoft/phi-2"
ADAPTER_PATH = "./final_phi_json_adapter" 
MAX_GENERATION_LENGTH = 1024 
GENERATION_TEMPERATURE = 0.45 # Lowered for stricter JSON adherence

# --- THEME DEFINITIONS (Unique Styling) ---
PALETTES = [
    {"name": "Tech Glacier", "bg": RGBColor(240, 248, 255), "text": RGBColor(50, 50, 50), "accent": RGBColor(0, 150, 200), "title": RGBColor(0, 51, 102), "font": "Calibri"},
    {"name": "Vibrant Sunset", "bg": RGBColor(255, 230, 200), "text": RGBColor(40, 40, 40), "accent": RGBColor(255, 120, 0), "title": RGBColor(160, 32, 240), "font": "Arial"},
    {"name": "Classic Dark", "bg": RGBColor(25, 25, 25), "text": RGBColor(240, 240, 240), "accent": RGBColor(0, 255, 127), "title": RGBColor(255, 255, 255), "font": "Verdana"}
]

class ThemeEngine:
    """Handles dynamic styling for the presentation."""
    def __init__(self):
        self.palette = random.choice(PALETTES)
        self.font_face = self.palette["font"]

    def apply_style(self, run, color_key="text", size=18, bold=False):
        run.font.name = self.font_face
        run.font.size = Pt(size)
        run.font.color.rgb = self.palette[color_key]
        run.font.bold = bold

    def set_slide_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.palette["bg"]

# --- QWEN IMAGE MATCHING SIMULATION ---
def find_best_image_match(image_concept, images_map):
    """ Placeholder for your Qwen embedding search."""
    # In a real service, this uses vector similarity search. Here, we simulate linking the ID to the path.
    for img_id, img_path in images_map.items():
        if img_id in image_concept: 
            return img_path
    return images_map.get("default", "extracted_images/default.png") # Use the default if available


class PPTX_Generator_Service:
    """
    A unified service to load the fine-tuned LLM and generate the PPTX file.
    """
    def __init__(self):
        self.model, self.tokenizer = self._load_fine_tuned_model()
        print("\n✨ PPTX Generation Service Ready (Phi-2 + QLoRA).")

    def _load_fine_tuned_model(self):
        """Loads the base Phi model and attaches QLoRA adapters."""
        print("Loading fine-tuned Phi model and adapters...")
        
        # 1. Quantization Configuration
        bnb_config = BitsAndBytesConfig(
            load_in_4bit=True, bnb_4bit_quant_type="nf4", 
            bnb_4bit_compute_dtype=torch.bfloat16, bnb_4bit_use_double_quant=True
        )

        # 2. Load Tokenizer and Set Padding
        tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
        if tokenizer.pad_token is None:
            tokenizer.pad_token = tokenizer.eos_token
        tokenizer.padding_side = "left"

        # 3. Load Base Model (4-bit)
        base_model = AutoModelForCausalLM.from_pretrained(
            MODEL_NAME, quantization_config=bnb_config, device_map="auto"
        )

        # 4. Load LoRA Adapters
        model = PeftModel.from_pretrained(base_model, ADAPTER_PATH) 
        model.eval()
        
        return model, tokenizer

    # --- INFERENCE: JSON Generation (Step 2) ---
    def generate_ppt_json(self, raw_text: str, available_images: list) -> dict or None:
        """Generates the structured JSON blueprint from text."""
        instruction = ("You are a Presentation Architect. Output ONLY the JSON object, strictly adhering to the schema.")
        image_list_str = str(available_images)
        input_content = f"---INPUT TEXT---\n{raw_text}\n\n---AVAILABLE IMAGES---\n{image_list_str}"
        prompt = f"### Instruction:\n{instruction}\n\n### Input:\n{input_content}\n\n### Output:\n"
        
        inputs = self.tokenizer(prompt, return_tensors="pt", padding=True, truncation=True)
        inputs = {k: v.to(self.model.device) for k, v in inputs.items()}
        
        with torch.no_grad():
            outputs = self.model.generate(
                **inputs, max_new_tokens=MAX_GENERATION_LENGTH, do_sample=True,      
                temperature=GENERATION_TEMPERATURE, top_p=0.9, num_beams=1,
                eos_token_id=self.tokenizer.eos_token_id 
            )
            
        generated_text = self.tokenizer.decode(outputs[0], skip_special_tokens=True)
        
        # --- ROBUST JSON EXTRACTION LOGIC (The crucial fix) ---
        try:
            print(f"   Raw Gen Length: {len(generated_text)} chars")
            json_start_index = generated_text.find("### Output:")
            if json_start_index != -1:
                json_string_candidate = generated_text[json_start_index:].split("### Output:")[-1].strip()
            else:
                # Fallback: Model might have skipped the marker, look for first brace
                json_string_candidate = generated_text.strip()
            
            # Aggressive Cleanup
            json_string_candidate = json_string_candidate.strip()
            if json_string_candidate.startswith('```'):
                json_string_candidate = json_string_candidate[json_string_candidate.find('\n')+1:]
            if json_string_candidate.endswith('```'):
                json_string_candidate = json_string_candidate[:-3]
            
            first_brace = json_string_candidate.find('{')
            last_brace = json_string_candidate.rfind('}')
            
            if first_brace == -1 or last_brace == -1 or last_brace < first_brace:
                raise ValueError("JSON braces not found or malformed.")

            json_string = json_string_candidate[first_brace : last_brace + 1].strip()
            
            return json.loads(json_string)
            
        except Exception as e:
            print(f"--- FAILED TO PARSE JSON --- Error: {e}")
            return None

    # --- ASSEMBLY: PPTX Building (Step 3) ---
    def create_presentation_file(self, raw_text: str, image_map: dict, output_dir: str = "."):
        """
        Main function to orchestrate JSON generation and PPTX assembly.
        """
        print("1. Generating JSON blueprint...")
        ppt_data = self.generate_ppt_json(raw_text, list(image_map.keys()))
        
        if not ppt_data:
            print("❌ Aborted: Could not generate a valid presentation structure.")
            return None
            
        print(f"   ✅ JSON Parsed. Title: {ppt_data.get('presentation_title')}")
        print(f"   ✅ Slides Found: {len(ppt_data.get('slides', []))}")
            
        print("2. Assembling PPTX file...")
        prs = Presentation()
        theme = ThemeEngine()
        
        # --- Build Slides ---
        # A. Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        theme.set_slide_background(slide)
        title = slide.shapes.title
        title.text = ppt_data.get("presentation_title", "Untitled Presentation")
        theme.apply_style(title.text_frame.paragraphs[0].add_run(), "title", size=48, bold=True)
        subtitle = slide.placeholders[1]
        subtitle.text = ppt_data.get("presentation_subtitle", "Generated Summary")
        theme.apply_style(subtitle.text_frame.paragraphs[0].add_run(), "text", size=24)

        # B. Content Slides
        slides_list = ppt_data.get("slides", [])
        if not slides_list:
            print("⚠️ Warning: JSON contained no slides. PPT will only have a title.")

        for slide_data in slides_list:
            image_reservation = slide_data.get("layout_type") == "TITLE_AND_CONTENT_WITH_IMAGE"
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            theme.set_slide_background(slide)

            # 1. Main Title & Sub-Heading
            title_width = Inches(9)
            # Add Title Box
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), title_width, Inches(0.8))
            p_title = title_box.text_frame.paragraphs[0]
            p_title.text = slide_data.get("main_title", "")
            theme.apply_style(p_title.add_run(), "title", size=32, bold=True)
            # Add Sub-heading Box
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), title_width, Inches(0.4))
            p_sub = sub_box.text_frame.paragraphs[0]
            p_sub.text = slide_data.get("sub_heading", "")
            theme.apply_style(p_sub.add_run(), "accent", size=18, bold=False)
            
            # 2. Bullet Points
            text_width = Inches(5.5) if image_reservation else Inches(9)
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), text_width, Inches(5.0))
            tf = body_box.text_frame
            tf.word_wrap = True
            
            for i, point in enumerate(slide_data.get("bullet_points", [])):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = point
                p.level = 0
                theme.apply_style(p.add_run(), "text", size=16)

            # 3. Image Placement
            if image_reservation and slide_data.get("image_concept"):
                image_path = find_best_image_match(slide_data["image_concept"], image_map)
                
                if os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(6.0), Inches(2.0), height=Inches(4.0))
                else:
                    print(f"Warning: Image file not found at {image_path}. Skipping placement.")
                
            # 4. Speaker Notes
            slide.notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")

        # 3. Save Final File
        output_filename = os.path.join(output_dir, f"PPT_Generated_{theme.palette['name'].replace(' ', '')}.pptx")
        prs.save(output_filename)
        
        print("\n" + "="*60)
        print(f"🎉 FINAL SERVICE OUTPUT: {output_filename}")
        print(f"   Theme Applied: {theme.palette['name']}")
        print("="*60)
        return output_filename

# --- Example Execution Block ---
if __name__ == "__main__":
    
    # 1. Simulate data preparation inputs
    SAMPLE_RAW_TEXT = (
        "The primary challenge in adapting AI to healthcare is ensuring compliance with HIPAA regulations, "
        "particularly concerning patient data privacy. LLMs often require data to be de-identified before training, "
        "a process which can lead to a loss of valuable clinical context. Figure 4.1 shows the two-stage pipeline. "
    )
    
    # 2. Simulate image file availability (Marker/Qwen output)
    # NOTE: You MUST ensure these image files exist in the 'extracted_images' directory for the PPTX part to work!
    IMAGE_MAPPING = {
        "IMAGE_RES1_F1": "extracted_images/health_pipeline.png",
        "IMAGE_RES1_F2": "extracted_images/hipaa_chart.png",
        "default": "extracted_images/default.png" 
    }
    
    # Simple check to create dummy files for testing the PPTX builder
    if not os.path.exists("extracted_images"):
        os.makedirs("extracted_images")
        print("Note: Created dummy directory 'extracted_images'. Please place actual PNG files here.")
        # Create minimal placeholder files (Actual image data is required, this is just to prevent file not found errors)
        try:
            from PIL import Image
            for path in IMAGE_MAPPING.values():
                if path.endswith('.png'):
                    img = Image.new('RGB', (400, 300), color = 'red' if 'default' in path else 'blue')
                    img.save(path)
        except ImportError:
            print("Install Pillow (pip install Pillow) to generate placeholder images.")

    # 3. Initialize and Run the Service
    try:
        service = PPTX_Generator_Service()
        service.create_presentation_file(SAMPLE_RAW_TEXT, IMAGE_MAPPING)
    except Exception as e:
        print(f"\n--- FATAL SERVICE ERROR ---")
        print(f"Ensure Phi-2 is downloaded and the adapter path ({ADAPTER_PATH}) is correct.")
        print(f"Error: {e}")